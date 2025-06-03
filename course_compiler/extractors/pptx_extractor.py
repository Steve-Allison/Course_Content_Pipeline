"""
PowerPoint (.pptx) file extractor implementation.
"""
from pathlib import Path
from typing import Dict, Any, List, Optional
from pptx import Presentation
from PIL import Image
import io
import os

from .base import BaseExtractor, ExtractionError
from ..logging_utils import get_logger

logger = get_logger(__name__)

class PptxExtractor(BaseExtractor):
    """Extracts content from PowerPoint (.pptx) files."""

    def __init__(self, file_path: Path, config: Optional[Dict[str, Any]] = None):
        super().__init__(file_path, config)
        self.presentation = None
        self.images_dir = self.config.get('images_dir', 'images')

    def load(self) -> None:
        """Load the PowerPoint presentation."""
        try:
            self.presentation = Presentation(self.file_path)
        except Exception as e:
            raise ExtractionError(f"Failed to load PowerPoint file: {e}") from e

    def extract(self) -> Dict[str, Any]:
        """
        Extract content from the PowerPoint file.

        Returns:
            Dictionary containing slides, notes, and metadata
        """
        if self.presentation is None:
            self.load()

        try:
            metadata = self.get_metadata()
            slides = []

            # Create output directory for images if it doesn't exist
            os.makedirs(self.images_dir, exist_ok=True)

            for i, slide in enumerate(self.presentation.slides):
                slide_data = self._extract_slide(slide, i)
                slides.append(slide_data)

            return {
                "metadata": metadata,
                "slides": slides,
                "total_slides": len(slides),
                "file_type": "pptx"
            }

        except Exception as e:
            logger.error(f"Error extracting content from {self.file_path}: {e}")
            raise ExtractionError(f"Failed to extract content: {e}") from e

    def _extract_slide(self, slide, slide_index: int) -> Dict[str, Any]:
        """Extract content from a single slide."""
        slide_data = {
            "slide_number": slide_index + 1,
            "shapes": [],
            "notes": self._extract_notes(slide)
        }

        # Extract title if available
        title_shapes = [s for s in slide.shapes if s.has_text_frame and s.text.strip()]
        if title_shapes:
            slide_data["title"] = title_shapes[0].text.strip()

        # Process all shapes in the slide
        for shape in slide.shapes:
            shape_data = self._process_shape(shape, slide_index)
            if shape_data:
                slide_data["shapes"].append(shape_data)

        return slide_data

    def _extract_notes(self, slide) -> Optional[str]:
        """Extract speaker notes from a slide."""
        try:
            if hasattr(slide, "has_notes_slide") and slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                return slide.notes_slide.notes_text_frame.text.strip()
        except Exception as e:
            logger.warning(f"Failed to extract notes from slide: {e}")
        return None

    def _process_shape(self, shape, slide_index: int) -> Optional[Dict[str, Any]]:
        """Process a single shape in the slide."""
        shape_data = {"type": shape.shape_type}

        try:
            if shape.has_text_frame and shape.text.strip():
                shape_data.update({
                    "type": "text",
                    "content": shape.text.strip(),
                    "size": (shape.width, shape.height),
                    "position": (shape.left, shape.top)
                })

            elif hasattr(shape, "image"):
                image_data = self._extract_image(shape, slide_index)
                if image_data:
                    shape_data.update({
                        "type": "image",
                        **image_data
                    })
                else:
                    return None

            else:
                # Skip shapes without content
                return None

            return shape_data

        except Exception as e:
            logger.warning(f"Failed to process shape: {e}")
            return None

    def _extract_image(self, shape, slide_index: int) -> Optional[Dict[str, Any]]:
        """Extract and save image from shape."""
        try:
            image = shape.image
            image_bytes = image.blob

            # Generate a unique filename for the image
            image_name = f"slide_{slide_index + 1}_img_{hash(image_bytes) & 0xFFFFFFFF}.png"
            image_path = os.path.join(self.images_dir, image_name)

            # Save the image
            with open(image_path, "wb") as f:
                f.write(image_bytes)

            # Get image dimensions
            with Image.open(io.BytesIO(image_bytes)) as img:
                width, height = img.size

            return {
                "path": image_path,
                "width": width,
                "height": height,
                "format": "png"
            }

        except Exception as e:
            logger.warning(f"Failed to extract image: {e}")
            return None
