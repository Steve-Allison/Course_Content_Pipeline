from pathlib import Path
from pptx import Presentation
from docx import Document
import fitz
import json
import os
from course_compiler.config import OUTPUT_ROOT
from .nlp_tagging import tag_text_components
from .ocr import extract_ocr_from_image
from .file_utils import sanitize_filename  # PATCH: Import sanitizer
import camelot
from PIL import Image

def get_file_metadata(file_path):
    p = Path(file_path)
    stat = p.stat()
    return {
        "file_name": p.name,
        "file_size": stat.st_size,
        "last_modified": stat.st_mtime,
        "created": stat.st_ctime,
        "absolute_path": str(p.resolve())
    }

def extract_pptx_hierarchical(pptx_path, out_dir):
    images_dir = Path(OUTPUT_ROOT) / "instructional_json" / "images"
    images_dir.mkdir(parents=True, exist_ok=True)
    prs = Presentation(pptx_path)
    slides = []
    for i, slide in enumerate(prs.slides):
        slide_data = {"slide_number": i + 1}
        title_shapes = [shape for shape in slide.shapes if shape.has_text_frame and shape.text_frame.text.strip()]
        slide_data["title"] = title_shapes[0].text.strip() if title_shapes else ""
        components = []
        # Speaker notes
        if hasattr(slide, "has_notes_slide") and slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_data["speaker_notes"] = notes_text
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                text = shape.text_frame.text.strip()
                tagged = tag_text_components(text)
                components.append(tagged)
            elif hasattr(shape, "image"):
                # PATCH: Sanitize shape name for file safety
                image_name = f"slide_{i+1}_{sanitize_filename(shape.name)}.png"
                image_path = images_dir / image_name
                with open(image_path, "wb") as imgf:
                    imgf.write(shape.image.blob)
                ocr_text = extract_ocr_from_image(image_path)
                image_metadata = get_image_metadata(image_path)
                components.append({
                    "type": "image",
                    "image_name": str(image_path.relative_to(OUTPUT_ROOT)),
                    "ocr_text": ocr_text,
                    "image_metadata": image_metadata
                })
            elif shape.shape_type == 19:  # Table
                table_data = []
                for row in shape.table.rows:
                    table_data.append([cell.text for cell in row.cells])
                components.append({
                    "type": "table",
                    "table": table_data
                })
            elif shape.has_chart:
                chart = shape.chart
                # Save chart as image (requires rendering, or use python-pptx-extract)
                # Or extract chart data:
                chart_data = []
                for series in chart.series:
                    series_data = [pt.value for pt in series.points]
                    chart_data.append({
                        "name": series.name,
                        "values": series_data
                    })
                components.append({
                    "type": "chart",
                    "chart_data": chart_data
                })
        slide_data["components"] = components
        slides.append(slide_data)
    return {
        "file_name": str(Path(pptx_path).name),
        "slides": slides,
        "file_metadata": get_file_metadata(pptx_path)
    }

def extract_docx_hierarchical(docx_path):
    doc = Document(docx_path)
    structure = []
    section = {"heading": None, "components": []}
    # Extract tables
    tables = []
    for table in doc.tables:
        table_data = []
        for row in table.rows:
            table_data.append([cell.text for cell in row.cells])
        tables.append(table_data)
    # Extract images
    images = []
    rels = doc.part.rels
    for rel in rels:
        rel_obj = rels[rel]
        if "image" in rel_obj.target_ref:
            images.append(rel_obj.target_ref)
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = para.style.name.lower()
        if "heading" in style:
            if section["components"]:
                structure.append(section)
            section = {"heading": text, "components": []}
        else:
            tagged = tag_text_components(text)
            section["components"].append(tagged)
    if section["components"]:
        structure.append(section)
    return {
        "file_name": str(Path(docx_path).name),
        "sections": structure,
        "tables": tables,
        "images": images,
        "file_metadata": get_file_metadata(docx_path)
    }

def extract_pdf_hierarchical(pdf_path, out_dir):
    images_dir = Path(OUTPUT_ROOT) / "instructional_json" / "images"
    images_dir.mkdir(parents=True, exist_ok=True)
    doc = fitz.open(pdf_path)
    pages = []
    for i, page in enumerate(doc):
        page_data = {"page_number": i+1, "components": []}
        text = page.get_text()
        if text.strip():
            tagged = tag_text_components(text)
            page_data["components"].append(tagged)
        images = page.get_images(full=True)
        for img_idx, img in enumerate(images):
            xref = img[0]
            pix = fitz.Pixmap(doc, xref)
            # PATCH: All names are numeric, but sanitize as a habit
            img_name = f"page_{i+1}_img_{img_idx+1}.png"
            img_path = images_dir / img_name
            pix.save(str(img_path))
            ocr_text = extract_ocr_from_image(img_path)
            image_metadata = get_image_metadata(img_path)
            page_data["components"].append({
                "type": "image",
                "image_name": str(img_path.relative_to(OUTPUT_ROOT)),
                "ocr_text": ocr_text,
                "image_metadata": image_metadata
            })
        # Basic table detection: look for lines/boxes (not robust, but a start)
        # This is a placeholder; for real table extraction, use a PDF table extraction library
        if "Table" in text or "table" in text:
            page_data["possible_table"] = True
        pages.append(page_data)
    return {
        "file_name": str(Path(pdf_path).name),
        "pages": pages,
        "file_metadata": get_file_metadata(pdf_path)
    }

def extract_tables_from_pdf(pdf_path):
    tables = []
    try:
        # flavor='stream' works better for tables without lines, 'lattice' for tables with lines
        camelot_tables = camelot.read_pdf(str(pdf_path), pages='all', flavor='stream')
        for table in camelot_tables:
            tables.append(table.df.values.tolist())  # as list of lists
    except Exception as e:
        print(f"Failed to extract tables from {pdf_path}: {e}")
    return tables

def save_instructional_json(data, output_dir, output_filename=None):
    # PATCH: Always sanitize output filename, either explicit or derived from data
    if output_filename is not None:
        out_path = Path(output_dir) / sanitize_filename(output_filename)
    else:
        out_path = Path(output_dir) / (sanitize_filename(data["file_name"]) + ".json")
    with open(out_path, "w") as f:
        json.dump(data, f, indent=2)

def get_image_metadata(image_path):
    try:
        with Image.open(image_path) as img:
            return {
                "format": img.format,
                "size": img.size,
                "mode": img.mode
            }
    except Exception as e:
        return {"error": str(e)}
